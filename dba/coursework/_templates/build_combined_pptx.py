#!/usr/bin/env python3
"""Combined Module 2 deck — Daniela's 9 slides preserved, Yasir's 7 appended,
plus appendices. Fully editable PPTX so either partner can revise.

Usage: python3 build_combined_pptx.py <her.pptx> <out.pptx>
"""
import sys, math
from pptx import Presentation
from pptx.util import Inches as In, Pt, Emu
from pptx.dml.color import RGBColor as C
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BLUE=C(0x08,0x1E,0x3F); GOLD=C(0xB6,0x86,0x2C); INK=C(0x14,0x17,0x1C)
PAPER=C(0xFC,0xFB,0xF8); MUTE=C(0x76,0x7D,0x86); RULE=C(0xDD,0xD7,0xCC)
SOFT=C(0xE8,0xED,0xF4); WARM=C(0xFA,0xF3,0xE4); WHITE=C(0xFF,0xFF,0xFF)
BODYC=C(0x45,0x4B,0x54); TEAL=C(0x1C,0x6B,0x63); RUST=C(0x8C,0x3A,0x1B)
STEEL=C(0x9A,0xA4,0xAE); TEALS=C(0xE3,0xEF,0xEC); RUSTS=C(0xF7,0xEA,0xE4)
DEEP=C(0x0C,0x25,0x47); LINE=C(0x1B,0x32,0x52); PALEB=C(0x8F,0xA3,0xBC)
SERIF="Georgia"; MONO="Consolas"
SEAL="/home/user/Profile/assets/images/fiu-seal.png"

prs = Presentation(sys.argv[1])
W, H = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[0]   # her template ships one layout, no placeholders
M = In(0.55)

def slide(dark=False):
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    bg.fill.solid(); bg.fill.fore_color.rgb = BLUE if dark else PAPER
    bg.line.fill.background(); bg.shadow.inherit = False
    return s

def tb(s, x, y, w, h, txt, size=14, color=BODYC, bold=False, italic=False,
       font=SERIF, align=PP_ALIGN.LEFT, space=0, caps=False, lead=1.15):
    box = s.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, line in enumerate(txt.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = lead
        r = p.add_run(); r.text = line.upper() if caps else line
        f = r.font; f.name = font; f.size = Pt(size); f.color.rgb = color
        f.bold = bold; f.italic = italic
        if space:
            from pptx.oxml.ns import qn
            r.font._rPr.set('spc', str(int(space*100)))
    return box

def rect(s, x, y, w, h, fill=None, line=None, lw=1.0, dash=False):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sh.adjustments[0] = 0.04
    if fill: sh.fill.solid(); sh.fill.fore_color.rgb = fill
    else: sh.fill.background()
    if line:
        sh.line.color.rgb = line; sh.line.width = Pt(lw)
        if dash:
            from pptx.enum.dml import MSO_LINE_DASH_STYLE
            sh.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    else: sh.line.fill.background()
    sh.shadow.inherit = False
    return sh

def bar(s, x, y, w, h, fill):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.fill.background(); sh.shadow.inherit = False
    return sh

def dot(s, cx, cy, r, fill):
    sh = s.shapes.add_shape(MSO_SHAPE.OVAL, cx-r, cy-r, 2*r, 2*r)
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.fill.background(); sh.shadow.inherit = False
    return sh

def crest(s, light=False):
    try:
        s.shapes.add_picture(SEAL, M, In(0.3), height=In(0.42))
        ox = M + In(0.55)
    except Exception:
        bar(s, M, In(0.32), In(0.03), In(0.36), GOLD); ox = M + In(0.14)
    tb(s, ox, In(0.30), In(6), In(0.2), "FLORIDA INTERNATIONAL UNIVERSITY",
       size=9.5, color=WHITE if light else BLUE, bold=True, font=MONO, space=1.9)
    tb(s, ox, In(0.50), In(7), In(0.2),
       "Chapman Graduate School of Business  ·  Doctor of Business Administration",
       size=10, color=PALEB if light else MUTE)

def chrome(s, n, light=False):
    fg = C(0x7E,0x93,0xA8) if light else MUTE
    bar(s, M, H-In(0.52), W-2*M, Emu(9525), LINE if light else RULE)
    tb(s, M, H-In(0.44), In(4), In(0.2), "Yasir A. Malik  ·  Cohort 8.14",
       size=8.5, color=fg, bold=True, font=MONO)
    tb(s, M+In(4), H-In(0.44), In(4.3), In(0.2),
       "GEB 7365  ·  Module 2  ·  27 August 2026", size=8.5, color=fg,
       font=MONO, align=PP_ALIGN.CENTER)
    tb(s, W-M-In(1), H-In(0.44), In(1), In(0.2), str(n), size=8.5, color=GOLD,
       bold=True, font=MONO, align=PP_ALIGN.RIGHT)

def eyebrow(s, y, txt, color=GOLD):
    tb(s, M, y, W-2*M, In(0.22), txt, size=9.5, color=color, bold=True,
       font=MONO, space=2.0, caps=True)

def title(s, y, txt, size=29, color=INK):
    tb(s, M, y, W-2*M, In(1.0), txt, size=size, color=color, bold=True, lead=1.12)

def sigmoid(s, x0, y0, w, h, color=BLUE, width=3.0):
    """Editable freeform S-curve: down, up, down."""
    def ease(t): return t*t*(3-2*t)
    def cy(f):
        if f < 0.28: return 0.52-0.42*ease(f/0.28)
        if f < 0.72: return 0.10+0.80*ease((f-0.28)/0.44)
        return 0.90-0.34*ease((f-0.72)/0.28)
    pts = [(x0+w*(i/60.0), y0+h-h*cy(i/60.0)*0.94) for i in range(61)]
    ff = s.shapes.build_freeform(pts[0][0], pts[0][1])
    ff.add_line_segments(pts[1:], close=False)
    sh = ff.convert_to_shape()
    sh.fill.background(); sh.line.color.rgb = color; sh.line.width = Pt(width)
    sh.shadow.inherit = False
    return sh

# ═══ Insert both presenter names on Daniela's title slide ═══
t0 = prs.slides[0]
tb(t0, M, H-In(1.35), W-2*M, In(0.3),
   "DANIELA GARCIA AGUIRRE   ·   YASIR A. MALIK",
   size=11, color=GOLD, bold=True, font=MONO, space=1.8)
tb(t0, M, H-In(1.05), W-2*M, In(0.3),
   "FIU Chapman Graduate School of Business  ·  DBA Cohort 8.14  ·  27 August 2026",
   size=10.5, color=MUTE)

N = len(prs.slides)   # her slides end here

# ═══ 10 · HANDOFF ═══
s = slide(dark=True); crest(s, light=True); N += 1
eyebrow(s, In(1.15), "Part two  ·  picking up where Daniela ended")
tb(s, M, In(1.42), W-2*M-In(0.4), In(0.7),
   "Distance can be priced. Can it be measured?",
   size=32, color=WHITE, bold=True, lead=1.1)
rect(s, M, In(2.25), W-2*M, In(0.78), DEEP, LINE)
tb(s, M+In(0.22), In(2.38), W-2*M-In(0.44), In(0.2), "DANIELA JUST SAID",
   size=9, color=C(0x7E,0x93,0xA8), bold=True, font=MONO, space=1.5)
tb(s, M+In(0.22), In(2.58), W-2*M-In(0.44), In(0.4),
   "“Who owns and runs the firm changes how that same distance is perceived and acted on.”",
   size=15, color=WHITE, italic=True)
rect(s, M, In(3.20), W-2*M, In(0.72), C(0x12,0x32,0x5C), GOLD, 1.4)
tb(s, M+In(0.22), In(3.32), W-2*M-In(0.44), In(0.2), "SO MY QUESTION",
   size=9, color=GOLD, bold=True, font=MONO, space=1.5)
tb(s, M+In(0.22), In(3.52), W-2*M-In(0.44), In(0.4),
   "When the field tried to measure any of this — did it measure what it thought it was measuring?",
   size=15, color=WHITE, italic=True)
bw = (W-2*M-In(0.2))/2
for i,(nm,yr,sub) in enumerate([("CONTRACTOR, KUNDU & HSU","2003","the shape of the relationship"),
                                ("MEZIAS","2002","the cost of being foreign")]):
    x = M+i*(bw+In(0.2))
    rect(s, x, In(4.15), bw, In(0.85), DEEP, LINE)
    tb(s, x, In(4.28), bw, In(0.2), nm, size=10, color=GOLD, bold=True,
       font=MONO, align=PP_ALIGN.CENTER, space=1.2)
    tb(s, x, In(4.50), bw, In(0.2), yr, size=12, color=C(0x5E,0x77,0x93),
       align=PP_ALIGN.CENTER)
    tb(s, x, In(4.72), bw, In(0.2), sub, size=12.5, color=WHITE, italic=True,
       align=PP_ALIGN.CENTER)
eyebrow(s, In(5.22), "Then the conclusion — all four papers, one diagnosis")
chrome(s, N, light=True)

# ═══ 11 · THE SIGMOID ═══
s = slide(); crest(s); N += 1
eyebrow(s, In(1.05), "Contractor, Kundu & Hsu (2003) · JIBS 34(1): 5–18 · after Figure 1")
title(s, In(1.28), "Four findings. One curve.")
cx, cy0, cw, ch = M+In(0.55), In(2.10), In(7.1), In(2.35)
for f0,f1,tint in [(0.05,0.26,C(0xF6,0xEB,0xE6)),(0.34,0.66,C(0xED,0xF3,0xF2)),(0.74,0.97,C(0xF2,0xEC,0xE2))]:
    bar(s, cx+int(cw*f0), cy0, int(cw*(f1-f0)), ch, tint)
bar(s, cx, cy0+ch, cw, Emu(12700), STEEL)
bar(s, cx, cy0, Emu(12700), ch, STEEL)
sigmoid(s, cx, cy0, cw, ch)
for f0,f1,lab,col in [(0.05,0.26,"1",RUST),(0.34,0.66,"2",TEAL),(0.74,0.97,"3",GOLD)]:
    tb(s, cx+int(cw*((f0+f1)/2))-In(0.2), cy0+In(0.08), In(0.4), In(0.3), lab,
       size=15, color=col, bold=True, align=PP_ALIGN.CENTER)
tb(s, cx, cy0+ch+In(0.08), cw, In(0.2), "DEGREE OF INTERNATIONALIZATION",
   size=8.5, color=MUTE, bold=True, font=MONO, align=PP_ALIGN.CENTER, space=1.2)
sx = cx+cw+In(0.35)
tb(s, sx, In(2.05), In(3.4), In(0.2), "SAMPLE ONLY…", size=9, color=MUTE,
   bold=True, font=MONO, space=1.4)
yy = In(2.32)
for lab,res,col in [("segment 1","you publish a NEGATIVE linear result",RUST),
                    ("segment 2","you publish a POSITIVE linear result",TEAL),
                    ("segments 1+2","you publish a U",TEAL),
                    ("segments 2+3","you publish an INVERTED U",GOLD)]:
    dot(s, sx+In(0.05), yy+In(0.08), In(0.045), col)
    tb(s, sx+In(0.18), yy, In(3.2), In(0.18), lab, size=9, color=col, bold=True,
       font=MONO, space=1.0)
    tb(s, sx+In(0.18), yy+In(0.19), In(3.2), In(0.4), res, size=12, color=BODYC)
    yy += In(0.62)
rect(s, M, In(4.72), W-2*M, In(0.62), SOFT)
tb(s, M+In(0.2), In(4.85), W-2*M-In(0.4), In(0.5),
   "Thirty years of contradiction, and none of it was about firms. The literature was sampling different segments of one sigmoid and reporting each as a finding about the whole.",
   size=13.5, color=INK, italic=True)
chrome(s, N)

# ═══ 12 · WHO OVER-EXPANDS ═══
s = slide(); crest(s); N += 1
eyebrow(s, In(1.05), "Contractor, Kundu & Hsu (2003) — the finding that reaches practice")
title(s, In(1.28), "Some firms walk into stage three more easily.")
cx, cy0, cw, ch = M+In(0.55), In(2.15), In(6.9), In(2.15)
bar(s, cx+int(cw*0.72), cy0, int(cw*0.28), ch, C(0xF6,0xEB,0xE6))
bar(s, cx, cy0+ch, cw, Emu(12700), STEEL)
bar(s, cx, cy0, Emu(12700), ch, STEEL)
sigmoid(s, cx, cy0, cw, ch)
def curve_pt(f, x0, y0, w, h):
    def ease(t): return t*t*(3-2*t)
    if f<0.28: v=0.52-0.42*ease(f/0.28)
    elif f<0.72: v=0.10+0.80*ease((f-0.28)/0.44)
    else: v=0.90-0.34*ease((f-0.72)/0.28)
    return x0+int(w*f), y0+h-int(h*v*0.94)
for f,col in [(0.86,RUST),(0.48,TEAL)]:
    px,py = curve_pt(f,cx,cy0,cw,ch)
    dot(s, px, py, In(0.075), col); dot(s, px, py, In(0.028), PAPER)
tb(s, cx+int(cw*0.72), cy0-In(0.26), int(cw*0.28), In(0.2),
   "STAGE 3 — THE EXPENSIVE ONE", size=8.5, color=RUST, bold=True, font=MONO,
   align=PP_ALIGN.CENTER, space=1.2)
tb(s, cx, cy0+ch+In(0.08), cw, In(0.2), "DEGREE OF INTERNATIONALIZATION",
   size=8.5, color=MUTE, bold=True, font=MONO, align=PP_ALIGN.CENTER, space=1.2)
sx = cx+cw+In(0.3)
dot(s, sx+In(0.06), In(2.28), In(0.055), RUST)
tb(s, sx+In(0.2), In(2.20), In(3.2), In(0.2), "KNOWLEDGE-BASED SERVICES",
   size=9, color=RUST, bold=True, font=MONO, space=1.1)
tb(s, sx+In(0.2), In(2.40), In(3.2), In(1.0),
   "Advertising · market research · securities · publishing. “More driven by intangible assets” with “a much lower fixed capital cost burden” — one more market is cheap to say yes to.",
   size=11.5, color=BODYC)
dot(s, sx+In(0.06), In(3.55), In(0.055), TEAL)
tb(s, sx+In(0.2), In(3.47), In(3.2), In(0.2), "CAPITAL-INTENSIVE",
   size=9, color=TEAL, bold=True, font=MONO, space=1.1)
tb(s, sx+In(0.2), In(3.67), In(3.2), In(0.9),
   "Airlines · hotels · construction · retail. Fixed assets make every additional market an expensive commitment, so they stop sooner.",
   size=11.5, color=BODYC)
rect(s, M, In(4.62), W-2*M, In(0.66), WARM, C(0xE4,0xD3,0xA8))
tb(s, M+In(0.2), In(4.76), W-2*M-In(0.4), In(0.5),
   "“Few companies possess the managerial tools (such as regression plotting of firms in a sector) that would tell them when they have over-internationalized.”",
   size=14, color=INK, italic=True)
tb(s, M, In(5.40), W-2*M, In(0.2),
   "→ “SHOULD WE INTERNATIONALISE?” IS THE WRONG QUESTION. “WHERE ARE WE ON THE CURVE?” IS THE RIGHT ONE.",
   size=8.5, color=BLUE, bold=True, font=MONO, space=1.0)
chrome(s, N)

# ═══ 13 · MEZIAS ═══
s = slide(); crest(s); N += 1
eyebrow(s, In(1.05), "Mezias (2002) · SMJ 23(3): 229–244 · the craft lesson")
title(s, In(1.28), "He did not fix the measure. He replaced it.")
lw = (W-2*M-In(0.4))/2
rect(s, M, In(2.05), lw, In(1.75), WHITE, RULE)
tb(s, M+In(0.2), In(2.20), lw-In(0.4), In(0.2), "WHAT PERFORMANCE MEASURES DO",
   size=9, color=MUTE, bold=True, font=MONO, space=1.2)
for i,(lab,col) in enumerate([("advantages",TEAL),("disadvantages",RUST),("transfer pricing",STEEL)]):
    ch2 = rect(s, M+In(0.2), In(2.50)+i*In(0.32), In(1.5), In(0.24), col)
    tb(s, M+In(0.2), In(2.53)+i*In(0.32), In(1.5), In(0.2), lab, size=8,
       color=WHITE, bold=True, font=MONO, align=PP_ALIGN.CENTER)
rect(s, M+In(2.15), In(2.72), lw-In(2.4), In(0.42), C(0xE9,0xE5,0xDC), STEEL)
tb(s, M+In(2.15), In(2.82), lw-In(2.4), In(0.2), "ONE NUMBER", size=9, color=INK,
   bold=True, font=MONO, align=PP_ALIGN.CENTER, space=1.2)
tb(s, M+In(2.15), In(3.28), lw-In(2.4), In(0.2), "nothing isolable", size=11.5,
   color=MUTE, italic=True, align=PP_ALIGN.CENTER)
rx = M+lw+In(0.4)
rect(s, rx, In(2.05), lw, In(1.75), WHITE, BLUE, 1.6)
tb(s, rx+In(0.2), In(2.20), lw-In(0.4), In(0.2), "WHAT LABOUR LAWSUITS DO",
   size=9, color=BLUE, bold=True, font=MONO, space=1.2)
rect(s, rx+In(0.2), In(2.50), In(1.5), In(0.24), RUST)
tb(s, rx+In(0.2), In(2.53), In(1.5), In(0.2), "disadvantages", size=8, color=WHITE,
   bold=True, font=MONO, align=PP_ALIGN.CENTER)
for i,lab in enumerate(["advantages","transfer pricing"]):
    rect(s, rx+In(0.2), In(2.82)+i*In(0.32), In(1.5), In(0.24), None, STEEL, 1.0, dash=True)
    tb(s, rx+In(0.2), In(2.85)+i*In(0.32), In(1.5), In(0.2), lab, size=8,
       color=STEEL, bold=True, font=MONO, align=PP_ALIGN.CENTER)
    bar(s, rx+In(0.2), In(2.94)+i*In(0.32), In(1.5), Emu(9525), STEEL)
rect(s, rx+In(2.15), In(2.72), lw-In(2.4), In(0.42), SOFT, BLUE)
tb(s, rx+In(2.15), In(2.82), lw-In(2.4), In(0.2), "ONE CONSTRUCT", size=9,
   color=BLUE, bold=True, font=MONO, align=PP_ALIGN.CENTER, space=1.2)
tb(s, rx+In(2.15), In(3.28), lw-In(2.4), In(0.2), "labour-related disadvantage only",
   size=11, color=BLUE, italic=True, align=PP_ALIGN.CENTER)
rect(s, M, In(3.95), W-2*M, In(0.52), SOFT)
tb(s, M+In(0.2), In(4.06), W-2*M-In(0.4), In(0.4),
   "“While most performance measures aggregate advantages and disadvantages, labor lawsuits only measure labor-related disadvantage.”",
   size=13.5, color=INK, italic=True)
bw = (W-2*M-In(0.3))/3
for i,(v,k) in enumerate([("486","foreign subsidiaries"),("486","matched US firms"),("$600k+","average jury award")]):
    x = M+i*(bw+In(0.15))
    rect(s, x, In(4.60), bw, In(0.72), WHITE, RULE)
    tb(s, x, In(4.70), bw, In(0.35), v, size=23, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    tb(s, x, In(5.08), bw, In(0.2), k, size=8.5, color=MUTE, bold=True,
       font=MONO, align=PP_ALIGN.CENTER, space=1.0)
tb(s, M, In(5.44), W-2*M, In(0.2),
   "MATCHED ON INDUSTRY AND CITY — SO THE ONLY REMAINING DIFFERENCE IS FOREIGNNESS",
   size=8.5, color=BLUE, bold=True, font=MONO, space=1.0)
chrome(s, N)

# ═══ 14 · CONCLUSION ═══
s = slide(); crest(s); N += 1
eyebrow(s, In(1.05), "Conclusion — all four")
title(s, In(1.28), "Daniela already named it. It is true of all four.")
rect(s, M, In(2.02), W-2*M, In(0.58), TEALS, TEAL)
tb(s, M+In(0.2), In(2.13), W-2*M-In(0.4), In(0.2), "HER SLIDE 7", size=9,
   color=TEAL, bold=True, font=MONO, space=1.4)
tb(s, M+In(0.2), In(2.31), W-2*M-In(0.4), In(0.25),
   "“Contradictions trace to unmeasured heterogeneity, not weak theory.”",
   size=14.5, color=INK, italic=True)
hx1, hx2 = M+In(1.55), M+In(5.7)
tb(s, hx1, In(2.82), In(3.5), In(0.2), "WHAT WAS MEASURED", size=9, color=MUTE,
   bold=True, font=MONO, space=1.4)
tb(s, hx2, In(2.82), In(4), In(0.2), "BELIEVED TO BE", size=9, color=MUTE,
   bold=True, font=MONO, space=1.4)
bar(s, M, In(3.04), W-2*M, Emu(9525), RULE)
yy = In(3.18)
for nm,got,thought,col in [("CONTRACTOR","segments of one curve","the shape of the relationship",BLUE),
                           ("MEZIAS","aggregate performance","the cost of foreignness",TEAL),
                           ("GHEMAWAT","market size","opportunity",GOLD),
                           ("ARREGLE","220 studies","a settled question",RUST)]:
    bar(s, M, yy, In(0.045), In(0.34), col)
    tb(s, M+In(0.16), yy+In(0.07), In(1.4), In(0.2), nm, size=9.5, color=col,
       bold=True, font=MONO, space=1.1)
    tb(s, hx1, yy+In(0.06), In(3.9), In(0.25), got, size=13.5, color=INK)
    tb(s, hx2-In(0.55), yy+In(0.07), In(0.45), In(0.2), "→", size=13, color=STEEL)
    tb(s, hx2, yy+In(0.06), In(4.2), In(0.25), thought, size=13.5, color=MUTE, italic=True)
    yy += In(0.44)
rect(s, M, In(5.02), W-2*M, In(0.72), WARM, C(0xE4,0xD3,0xA8))
tb(s, M+In(0.2), In(5.15), W-2*M-In(0.4), In(0.5),
   "Not one of them was fixed by more data. Every one made progress by changing WHAT was measured — which is why this is a methods problem, not a theory problem.",
   size=15, color=INK, italic=True)
chrome(s, N)

# ═══ 15 · CLOSE ═══
s = slide(dark=True); crest(s, light=True); N += 1
eyebrow(s, In(1.12), "One last thing")
tb(s, M, In(1.38), W-2*M, In(0.6), "Eighteen years. Same question.",
   size=30, color=WHITE, bold=True)
tx0, tx1, ty = M+In(1.35), W-M-In(1.35), In(3.05)
bar(s, tx0, ty, tx1-tx0, Emu(25400), LINE)
for frac,yr,pap,ven in [(0.0,"2003","Contractor, Kundu & Hsu","JIBS 34(1)"),
                        (1.0,"2021","Arregle, … Kundu, …","JIBS 52")]:
    x = tx0+int((tx1-tx0)*frac)
    dot(s, x, ty, In(0.1), GOLD)
    tb(s, x-In(1.4), ty-In(0.85), In(2.8), In(0.4), yr, size=25, color=WHITE,
       bold=True, align=PP_ALIGN.CENTER)
    tb(s, x-In(1.4), ty-In(0.42), In(2.8), In(0.2), pap, size=9, color=GOLD,
       bold=True, font=MONO, align=PP_ALIGN.CENTER, space=1.0)
    tb(s, x-In(1.4), ty+In(0.2), In(2.8), In(0.2), ven, size=11.5,
       color=C(0x7E,0x93,0xA8), align=PP_ALIGN.CENTER)
mid = (tx0+tx1)//2
tb(s, mid-In(1.6), ty-In(1.05), In(3.2), In(0.2), "SUMIT K. KUNDU  ·  FIU",
   size=11, color=GOLD, bold=True, font=MONO, align=PP_ALIGN.CENTER, space=1.8)
tb(s, mid-In(1.6), ty+In(0.2), In(3.2), In(0.2), "author on both", size=12.5,
   color=PALEB, italic=True, align=PP_ALIGN.CENTER)
rect(s, M, In(3.72), W-2*M, In(0.6), DEEP, GOLD, 1.3)
tb(s, M+In(0.2), In(3.85), W-2*M-In(0.4), In(0.45),
   "The same question about multinationality and performance, asked of a different population eighteen years later — and Daniela's numbers show it is still the theme the field has looked at least.",
   size=14, color=WHITE, italic=True)
yy = In(4.50)
for d,k,v in [("3 SEP","ENTRY MODE","If stage one is negative, mode choice is a decision about how much of that cost to absorb at once."),
              ("24 SEP","INSTITUTIONS","CAGE is a pre-Hofstede way of saying what that session says with formal distance measures.")]:
    bar(s, M, yy, In(0.04), In(0.48), GOLD)
    tb(s, M+In(0.18), yy, In(0.9), In(0.2), d, size=9, color=GOLD, bold=True,
       font=MONO, space=1.3)
    tb(s, M+In(1.0), yy, In(2.0), In(0.2), k, size=9, color=PALEB, bold=True,
       font=MONO, space=1.3)
    tb(s, M+In(0.18), yy+In(0.2), W-2*M-In(0.4), In(0.3), v, size=12.5,
       color=C(0xB8,0xC9,0xD9))
    yy += In(0.6)
chrome(s, N, light=True)

# ═════════════ APPENDICES ═════════════

# ═══ A1 · BACKUP — his own research ═══
s = slide(); crest(s); N += 1
eyebrow(s, In(1.05), "Appendix A  ·  backup — only if asked")
title(s, In(1.28), "The same error, in my own study.", size=28)
rect(s, M, In(2.02), W-2*M, In(0.56), SOFT)
tb(s, M+In(0.2), In(2.14), W-2*M-In(0.4), In(0.4),
   "Tversky & Kahneman (1974) — a wheel of fortune gave subjects 10 or 65 before they estimated the share of African countries in the UN. Median answers: 25 and 45.",
   size=13.5, color=BODYC)
lw = (W-2*M-In(0.3))/2
rect(s, M, In(2.78), lw, In(1.85), WHITE, RULE); bar(s, M, In(2.78), lw, In(0.06), BLUE)
tb(s, M+In(0.2), In(2.98), lw-In(0.4), In(0.2), "WHAT MY MODEL MEASURES",
   size=9, color=BLUE, bold=True, font=MONO, space=1.2)
tb(s, M+In(0.2), In(3.20), lw-In(0.4), In(1.3),
   "Eight organisational interventions — training, rotation, analytical tools, structured processes, feedback, independent review, regulatory guidance, incentives — against anchoring bias, with expertise and confidence as moderators.",
   size=12, color=BODYC)
x2 = M+lw+In(0.3)
rect(s, x2, In(2.78), lw, In(1.85), WHITE, RUST, 1.6); bar(s, x2, In(2.78), lw, In(0.06), RUST)
tb(s, x2+In(0.2), In(2.98), lw-In(0.4), In(0.2), "WHAT IT ACTUALLY MEASURES",
   size=9, color=RUST, bold=True, font=MONO, space=1.2)
tb(s, x2+In(0.2), In(3.20), lw-In(0.4), In(1.0),
   "Anchoring bias is captured by self-report, reverse-coded. It asks auditors to report how far their judgment was driven by a reference point — the one thing the bias reliably prevents them from noticing.",
   size=12, color=BODYC)
tb(s, x2+In(0.2), In(4.24), lw-In(0.4), In(0.3),
   "Perceived judgment discipline, described as anchoring reduction.",
   size=11.5, color=RUST, italic=True)
rect(s, M, In(4.78), W-2*M, In(0.6), WARM, C(0xE4,0xD3,0xA8))
tb(s, M+In(0.2), In(4.90), W-2*M-In(0.4), In(0.45),
   "Contractor measured segments of a curve. Mezias's field measured aggregate performance. I measured self-reported bias resistance. Three instances of the same error, and one of them is mine.",
   size=14, color=INK, italic=True)
tb(s, M, In(5.48), W-2*M, In(0.2),
   "→ TVERSKY & KAHNEMAN ALREADY WARNED ME: “PAYOFFS FOR ACCURACY DID NOT REDUCE THE ANCHORING EFFECT.”",
   size=8.5, color=BLUE, bold=True, font=MONO, space=0.9)
chrome(s, N)

# ═══ A2 · VERIFICATION RECORD ═══
s = slide(); crest(s); N += 1
eyebrow(s, In(1.05), "Appendix B  ·  every quote and figure, with its page")
title(s, In(1.28), "Verification record.", size=28)
cols = [M+In(0.1), M+In(0.85), M+In(6.4), M+In(9.5)]
for i,h in enumerate(["SLIDE","CLAIM","SOURCE","CHECKED"]):
    tb(s, cols[i], In(2.02), In(3.0), In(0.2), h, size=8.5, color=MUTE, bold=True,
       font=MONO, space=1.2)
bar(s, M, In(2.22), W-2*M, Emu(9525), RULE)
yy = In(2.34)
for n,claim,src in [("11","Sigmoid, three stages; four findings are segments of one curve","Contractor p.7, Fig. 1"),
                    ("11","“Depending on which part of Figure 1 we examined…”","Contractor p.7"),
                    ("12","Knowledge-based vs capital-intensive; lower fixed capital burden","Contractor p.9"),
                    ("12","“Few companies possess the managerial tools…”","Contractor p.11"),
                    ("13","486 foreign subsidiaries vs 486 US, same industries and cities","Mezias p.229–231"),
                    ("13","“labor lawsuits only measure labor-related disadvantage”","Mezias p.231"),
                    ("13","78.9% / 70.0% / 58.4%; award > $600,000; fees ~$100,000","Mezias p.232"),
                    ("14","“varied and at times incompatible findings”","Arregle, abstract"),
                    ("15","Performance is 12% of the 220 studies","Arregle p.6"),
                    ("15","Kundu authored both the 2003 and 2021 papers","Both title pages"),
                    ("A1","Wheel of fortune: 10 → 25, 65 → 45; payoffs did not reduce it","Tversky & Kahneman p.1128")]:
    tb(s, cols[0], yy, In(0.7), In(0.2), n, size=9, color=BLUE, bold=True, font=MONO)
    tb(s, cols[1], yy, In(5.4), In(0.22), claim, size=11.5, color=INK)
    tb(s, cols[2], yy, In(3.0), In(0.22), src, size=11, color=MUTE)
    tb(s, cols[3], yy, In(1.2), In(0.2), "EXACT", size=8, color=TEAL, bold=True, font=MONO, space=1.0)
    bar(s, M, yy+In(0.22), W-2*M, Emu(6350), C(0xEF,0xEB,0xE3))
    yy += In(0.275)
rect(s, M, In(5.42), W-2*M, In(0.4), SOFT)
tb(s, M+In(0.2), In(5.50), W-2*M-In(0.4), In(0.25),
   "All verified against the source PDFs, 26 August 2026. Nothing is paraphrased and presented as a quote.",
   size=12, color=INK, italic=True)
chrome(s, N)

# ═══ A3 · REFERENCES ═══
s = slide(); crest(s); N += 1
eyebrow(s, In(1.05), "Appendix C")
title(s, In(1.28), "References.", size=28)
yy = In(2.10)
for who,ref in [("DANIELA","Ghemawat, P. (2001). Distance still matters: The hard reality of global expansion. Harvard Business Review, 79(8), 137–147."),
                ("DANIELA","Arregle, J.-L., Chirico, F., Kano, L., Kundu, S. K., Majocchi, A., & Schulze, W. S. (2021). Family firm internationalization: Past research and an agenda for the future. Journal of International Business Studies, 52, 1159–1198."),
                ("YASIR","Contractor, F. J., Kundu, S. K., & Hsu, C.-C. (2003). A three-stage theory of international expansion: The link between multinationality and performance in the service sector. Journal of International Business Studies, 34(1), 5–18."),
                ("YASIR","Mezias, J. M. (2002). Identifying liabilities of foreignness and strategies to minimize their effects: The case of labor lawsuit judgments in the United States. Strategic Management Journal, 23(3), 229–244."),
                ("APPENDIX A","Tversky, A., & Kahneman, D. (1974). Judgment under uncertainty: Heuristics and biases. Science, 185(4157), 1124–1131.")]:
    col = TEAL if who=="DANIELA" else (BLUE if who=="YASIR" else GOLD)
    bar(s, M, yy, In(0.04), In(0.5), col)
    tb(s, M+In(0.18), yy, In(1.3), In(0.2), who, size=8.5, color=col, bold=True,
       font=MONO, space=1.2)
    tb(s, M+In(1.7), yy-In(0.02), W-2*M-In(1.9), In(0.5), ref, size=12, color=BODYC)
    yy += In(0.62)
rect(s, M, In(5.20), W-2*M, In(0.6), SOFT)
tb(s, M+In(0.2), In(5.32), W-2*M-In(0.4), In(0.45),
   "Presented by Daniela Garcia Aguirre and Yasir A. Malik · GEB 7365 International Business Theory and Practice · Prof. William Newburry · FIU Chapman Graduate School of Business, DBA Cohort 8.14 · 27 August 2026",
   size=11.5, color=INK, italic=True)
chrome(s, N)

prs.save(sys.argv[2])
print(f"saved {sys.argv[2]} · {len(prs.slides)} slides")
