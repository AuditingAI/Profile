#!/usr/bin/env python3
"""Editable PPTX twin of RESEARCH_STATUS_DECK.pdf. Native shapes only."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

NAVY=RGBColor(0x08,0x1E,0x3F); GOLD=RGBColor(0xB6,0x86,0x2C); PAPER=RGBColor(0xFC,0xFB,0xF8)
INK=RGBColor(0x14,0x17,0x1C); MUTED=RGBColor(0x76,0x7D,0x86); RULE=RGBColor(0xD8,0xD4,0xCB)
BLUE=RGBColor(0x1F,0x4E,0x79); RED=RGBColor(0xA5,0x32,0x2F); GREEN=RGBColor(0x21,0x69,0x5A)
WHITE=RGBColor(0xFF,0xFF,0xFF); CARD=RGBColor(0xF2,0xF0,0xEC)

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
SW,SH=13.333,7.5; M=0.75
BLANK=prs.slide_layouts[6]

def slide(bg=PAPER):
    s=prs.slides.add_slide(BLANK)
    r=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,prs.slide_width,prs.slide_height)
    r.fill.solid(); r.fill.fore_color.rgb=bg; r.line.fill.background(); r.shadow.inherit=False
    return s

def tb(s,x,y,w,h,txt,size=10,color=INK,bold=False,italic=False,font="Georgia",
       align=PP_ALIGN.LEFT,space=2,lead=1.15):
    box=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h))
    tf=box.text_frame; tf.word_wrap=True; tf.margin_left=0; tf.margin_right=0
    tf.margin_top=0; tf.margin_bottom=0
    for i,line in enumerate(txt.split("\n")):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=align; p.space_after=Pt(space); p.line_spacing=lead
        r=p.add_run(); r.text=line
        f=r.font; f.name=font; f.size=Pt(size); f.color.rgb=color; f.bold=bold; f.italic=italic
    return box

def rect(s,x,y,w,h,fill=None,line=None,lw=1.0):
    sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(x),Inches(y),Inches(w),Inches(h))
    if fill: sh.fill.solid(); sh.fill.fore_color.rgb=fill
    else: sh.fill.background()
    if line: sh.line.color.rgb=line; sh.line.width=Pt(lw)
    else: sh.line.fill.background()
    sh.shadow.inherit=False
    return sh

def chrome(s,eyebrow,title,num,sub=None):
    rect(s,0,0,SW,0.055,fill=NAVY)
    tb(s,M,0.42,6,0.2,"FLORIDA INTERNATIONAL UNIVERSITY",7,NAVY,True,font="Consolas")
    tb(s,M,0.58,7,0.2,"Chapman Graduate School of Business  ·  Doctor of Business Administration",7,MUTED)
    tb(s,M,1.05,8,0.2,eyebrow,7.5,GOLD,True,font="Consolas")
    tb(s,M,1.32,11.5,0.5,title,22,NAVY,True)
    if sub: tb(s,M,1.82,11.5,0.3,sub,10.5,MUTED,italic=True)
    rect(s,M,SH-0.62,SW-2*M,0.008,fill=RULE)
    tb(s,M,SH-0.48,6,0.2,"Yasir A. Malik  ·  DBA Cohort 8.14",7,MUTED,font="Consolas")
    tb(s,SW-M-1,SH-0.48,1,0.2,num,7,MUTED,font="Consolas",align=PP_ALIGN.RIGHT)

# 1 title
s=slide(NAVY)
rect(s,M,2.05,0.06,0.65,fill=GOLD)
tb(s,M+0.22,2.12,7,0.2,"FLORIDA INTERNATIONAL UNIVERSITY",8,WHITE,True,font="Consolas")
tb(s,M+0.22,2.34,8,0.2,"Chapman Graduate School of Business  ·  DBA Cohort 8.14",8.5,RGBColor(0xA8,0xB8,0xCF))
tb(s,M,3.25,11,0.7,"Where the research stands,",38,WHITE,True)
tb(s,M,3.85,11,0.7,"and what unlocks the next step.",38,GOLD,True)
tb(s,M,4.75,8.4,1.0,"A status read on the anchoring-bias programme: the model as built, the constraint that "
   "stopped it, the AI extension, and the five approvals that gate everything downstream.",12,RGBColor(0xB8,0xC6,0xDB))
rect(s,M,5.75,4.4,0.02,fill=GOLD)
tb(s,M,6.0,6,0.3,"Yasir A. Malik",14,WHITE,True)
tb(s,M,6.32,6,0.2,"28 August 2026",8,RGBColor(0x8C,0x9E,0xB8),font="Consolas")

# 2 model
s=slide(); chrome(s,"THE MODEL AS BUILT","Eleven constructs, sixteen hypotheses.","2",
    "Eight organisational interventions, two mediators, one outcome. 55 Likert items.")
tb(s,M,2.32,2.2,0.2,"EIGHT INTERVENTIONS",6.5,GOLD,True,font="Consolas")
IV=[("TA","Training & Awareness"),("RA","Rotation of Auditors"),("AT","Analytical Tools"),
    ("SAP","Structured Processes"),("FR","Feedback & Reflection"),("IR","Independent Reviews"),
    ("RPG","Regulatory Guidance"),("PMI","Metrics & Incentives")]
for i,(k,lab) in enumerate(IV):
    y=2.55+i*0.50
    rect(s,M,y,2.15,0.44,fill=WHITE,line=RULE,lw=0.75)
    rect(s,M,y,0.035,0.44,fill=BLUE)
    tb(s,M+0.13,y+0.05,1.9,0.15,k,6.5,BLUE,True,font="Consolas")
    tb(s,M+0.13,y+0.20,1.95,0.18,lab,7.5,INK)
tb(s,4.42,2.32,3.0,0.2,"TWO MEDIATORS",6.5,GOLD,True,font="Consolas")
for i,(k,t,d) in enumerate([("AJQ","Auditor Judgment Quality","Cognitive — careful, objective evaluation of independent evidence"),
                            ("APR","Audit Process Rigor","Procedural — thorough, consistent, disciplined execution")]):
    y=2.72+i*1.45
    rect(s,4.42,y,3.15,1.18,fill=WHITE,line=BLUE,lw=1.1)
    tb(s,4.60,y+0.14,2,0.18,k,8,BLUE,True,font="Consolas")
    tb(s,4.60,y+0.38,2.8,0.25,t,11,INK,True)
    tb(s,4.60,y+0.68,2.8,0.5,d,8,MUTED)
tb(s,8.85,2.32,2.0,0.2,"OUTCOME",6.5,GOLD,True,font="Consolas")
rect(s,8.85,2.72,3.7,1.45,fill=NAVY,line=GOLD,lw=1.4)
tb(s,9.05,2.88,2,0.18,"RAB",8.5,GOLD,True,font="Consolas")
tb(s,9.05,3.12,3.3,0.5,"Reduction in\nAnchoring Bias",14,WHITE,True)
tb(s,9.05,3.72,3.3,0.4,"Final judgments driven by current evidence, not initial reference points.",8,RGBColor(0xB0,0xC0,0xD8))
rect(s,8.85,4.35,3.7,1.15,fill=RGBColor(0xFA,0xF0,0xEF),line=RED,lw=1.0)
tb(s,9.05,4.50,3.3,0.18,"THE MEASUREMENT PROBLEM",6.5,RED,True,font="Consolas")
tb(s,9.05,4.72,3.35,0.7,"RAB is captured by self-report. It asks auditors to report how far judgment "
   "was driven by a reference point — the one thing the bias reliably prevents them from noticing.",8,INK)
tb(s,M,6.6,11,0.2,"8 interventions × 2 mediators + direct paths = 16 hypotheses.   Each construct: 5 items, one reverse-coded.",7,MUTED,font="Consolas")

# 3 funnel
s=slide(); chrome(s,"WHAT ACTUALLY HAPPENED","The instrument worked. The population did not exist.","3",
    "Applying the eligibility criteria to a commercial research panel, July 2026.")
for i,(n,lab,w,col) in enumerate([("334,976","panel members screened",10.6,BLUE),
                                   ("~20","matched the eligibility criteria",5.1,BLUE),
                                   ("23","raw responses recorded (Qualtrics, organic outreach)",3.2,GOLD),
                                   ("4","survived screening",1.7,RED)]):
    y=2.45+i*0.86; x=M+(11.83-w)/2
    rect(s,x,y,w,0.70,fill=WHITE,line=col,lw=1.2)
    tb(s,x+0.18,y+0.16,1.6,0.4,n,21,col,True)
    tb(s,x+0.18+1.55,y+0.26,w-1.9,0.3,lab,9.5,MUTED)
rect(s,M,5.95,SW-2*M,0.95,fill=NAVY,line=GOLD,lw=1.2)
tb(s,M+0.22,6.10,7,0.2,"PREVALENCE ≈ 6 PER 100,000",7.5,GOLD,True,font="Consolas")
tb(s,M+0.22,6.36,11,0.25,"A survey at this prevalence needed a sampling frame of roughly 9.6 million. The panel held 3.5% of that.",9.5,RGBColor(0xE0,0xE8,0xF2))
tb(s,M+0.22,6.62,11,0.2,"Source: dba/03_Data/EXCLUSION_LOG_2026-07-22.md  ·  dba/RISK_QUANT/feasibility.py reproduces this end to end.",6.5,RGBColor(0x8C,0x9E,0xB8),font="Consolas")

# 4 chain
s=slide(); chrome(s,"THE EXTENSION","Three links. Only the first is a cognitive bias.","4",
    "What happens when the anchor stops being a workpaper and becomes a machine.")
L=[("L1","Automated anchoring","The anchor is system-generated — continuous, and arriving BEFORE the reviewer forms a view",
    "Automation bias","HOW MUCH · magnitude","Survey / experiment","Human cognitive bias",GREEN),
   ("L2","Sycophantic confirmation","The system AGREES with a position the auditor already stated, rather than challenging it",
    "Sycophancy","HOW · mechanism","Interviews — phenomenology","Model behaviour, not a bias",GOLD),
   ("L3","Recursive epistemic drift","Successive models reprocess earlier AI-influenced work; the evidentiary basis thins",
    "Model collapse","HOW, OVER TIME · process","Longitudinal / archival","Property of a system of models",RED)]
bw=(11.83-2*0.30)/3
for i,(k,t,d,ph,qt,me,note,col) in enumerate(L):
    x=M+i*(bw+0.30)
    rect(s,x,2.35,bw,3.15,fill=WHITE,line=col,lw=1.3)
    rect(s,x,2.35,bw,0.34,fill=col)
    tb(s,x+0.16,2.43,1,0.2,k,9.5,WHITE,True,font="Consolas")
    tb(s,x+0.16,2.82,bw-0.32,0.3,t,11.5,NAVY,True)
    tb(s,x+0.16,3.15,bw-0.32,0.6,d,8.2,INK)
    rect(s,x+0.16,3.86,bw-0.32,0.008,fill=RULE)
    for j,(lab,val,c,b) in enumerate([("PHENOMENON",ph,col,True),("QUESTION TYPE",qt,INK,False),("METHOD THAT FITS",me,INK,False)]):
        yy=4.00+j*0.44
        tb(s,x+0.16,yy,2,0.14,lab,5.8,MUTED,True,font="Consolas")
        tb(s,x+0.16,yy+0.15,bw-0.32,0.22,val,9,c,b,font="Consolas" if j==1 else "Georgia")
    tb(s,x+0.16,5.20,bw-0.32,0.2,note,7.6,MUTED,italic=True)
rect(s,M,5.66,SW-2*M,0.72,fill=RGBColor(0xF7,0xF4,0xEB),line=GOLD,lw=1.0)
tb(s,M+0.22,5.80,9,0.2,"THE PRECISION THAT MATTERS WHEN A PROFESSOR PUSHES",6.8,GOLD,True,font="Consolas")
tb(s,M+0.22,6.02,11.4,0.3,'Only L1 involves a human cognitive bias. L2 is model behaviour interacting with human judgment. '
   'L3 is a property of a system of models. Calling all three "biases" loses the argument.',9,INK)
tb(s,M,6.55,11,0.2,"L3 is deliberately NOT designed — it needs longitudinal access that does not exist. Stated as future work, not claimed.",6.8,MUTED,font="Consolas")

# 5 courses
s=slide(); chrome(s,"HOW THIS TERM FEEDS THE DISSERTATION","Two courses, two arms, no double-submission.","5",
    "Each course builds a real piece. Neither collects a single response.")
for i,(hdr,who,claim,body,bul,col) in enumerate([
    ("GEB 7365 · INTERNATIONAL BUSINESS","Prof. William Newburry","Feasibility as a design parameter",
     "Requires a data-collection plan and explicitly NO collected data — the first time the recruitment constraint can be worked on without being blocked by it.",
     ["Study model + 2–3 hypotheses","Cross-national feasibility argument","Presentation 19 Sep · Paper 9 Oct"],BLUE),
    ("GEB 7911 · QUALITATIVE METHODS","Dr. Gonzalez","The L2 phenomenology arm",
     "The research proposal is 50% of that grade and IS the L2 design — protocol, sampling, coding plan, trustworthiness, audit trail.",
     ["Proposal presentations 6 Oct","Learning memos, weekly","Method: phenomenology, her recommendation"],GOLD)]):
    x=M+i*(5.85+0.13)
    rect(s,x,2.35,5.85,2.75,fill=WHITE,line=col,lw=1.2)
    rect(s,x,2.35,5.85,0.055,fill=col)
    tb(s,x+0.22,2.52,5.4,0.18,hdr,7,col,True,font="Consolas")
    tb(s,x+0.22,2.72,5.4,0.18,who,7.6,MUTED)
    tb(s,x+0.22,2.98,5.4,0.3,claim,13,NAVY,True)
    tb(s,x+0.22,3.36,5.45,0.7,body,8.6,INK)
    for j,b in enumerate(bul):
        tb(s,x+0.22,4.12+j*0.28,0.2,0.2,"›",8,col,True,font="Consolas")
        tb(s,x+0.42,4.12+j*0.28,5.2,0.22,b,8.5,INK)
rect(s,M,5.30,SW-2*M,1.05,fill=RGBColor(0xFA,0xF0,0xEF),line=RED,lw=1.1)
tb(s,M+0.22,5.46,9,0.2,"WHAT NEITHER COURSE DOES — SAID PLAINLY",7,RED,True,font="Consolas")
tb(s,M+0.22,5.70,11.4,0.6,"Neither collects a response, so the achieved n is still four. Neither rebuilds Chapters 2 and 3 — "
   "those are still owed and Chapters 4–6 are overdue since 28 July. A course grade is not committee approval. "
   "And the 7365 topic STUDIES the feasibility problem; it does not solve it.",8.8,INK)

# 6 gates
s=slide(); chrome(s,"WHOSE APPROVAL UNLOCKS WHAT","Five gates. Four are open and one was never asked.","6",
    'The honest answer to "which approval do I need to level up."')
G=[("DR. JUAN REY","Advisor · chair","Chapters 4–6 · the Ch 2–3 rebuild · whether the survey arm relaunches · AI-use approval",
    "Ch 4–6 OVERDUE since 28 Jul.\nAI approval NEVER REQUESTED",RED,"THE BINDING ONE"),
   ("FIU IRB OFFICE","Office of Research Integrity","Whether interviews are a modification or a new protocol · approval before ANY participant",
    "3 items unsubmitted · 1 never confirmed filed",RED,"BLOCKS ALL L2 DATA"),
   ("DR. GONZALEZ","GEB 7911","The qualitative proposal — 50% of that grade · whether phenomenology fits the question",
    "Protocol drafted. She has never seen it",GOLD,"ASK EARLY"),
   ("PROF. NEWBURRY","GEB 7365","Project topic approval before extensive work",
    "Zoom Tue 1 Sep — this one is in motion",GREEN,"IN MOTION"),
   ("THE COMMITTEE","Dissertation committee","Proposal defence · AI-use review required by FIU UGS §3.2",
    "Composition [VERIFY] · UGS §3.2 review never requested",MUTED,"DOWNSTREAM")]
for i,(name,role,gates,status,col,tag) in enumerate(G):
    y=2.32+i*0.83
    rect(s,M,y,SW-2*M,0.76,fill=WHITE,line=RULE,lw=0.7)
    rect(s,M,y,0.05,0.76,fill=col)
    tb(s,M+0.22,y+0.09,2.4,0.18,name,8,NAVY,True,font="Consolas")
    tb(s,M+0.22,y+0.27,2.4,0.18,role,7.2,MUTED)
    tb(s,M+0.22,y+0.52,2.4,0.16,tag,5.8,col,True,font="Consolas")
    tb(s,3.35,y+0.09,0.8,0.14,"GATES",5.8,MUTED,True,font="Consolas")
    tb(s,3.35,y+0.26,5.3,0.45,gates,8.2,INK)
    tb(s,8.95,y+0.09,0.8,0.14,"STATUS",5.8,MUTED,True,font="Consolas")
    tb(s,8.95,y+0.26,3.6,0.45,status,8.2,col,True)
tb(s,M,6.62,11.5,0.2,"Two of the open items cost nothing: verifying the Topaz record, and asking the IRB office one question. A question is not a submission.",6.8,MUTED,font="Consolas")

# 7 model improvement
s=slide(); chrome(s,"IMPROVING THE MODEL","Three known weaknesses, and what each one actually needs.","7",
    "Written down because a model whose weaknesses are named is defensible; one whose are not is fragile.")
Wk=[("1","Construct 11 measures the wrong thing",
     "RAB — the outcome — is self-report, reverse-coded. It asks auditors to report how far their judgment was driven by an anchor. That is precisely what anchoring prevents them from noticing.",
     "Either rename it honestly to PERCEIVED JUDGMENT DISCIPLINE, or replace it with a behavioural measure: a task with a planted anchor and an observable adjustment.",
     "This is the single most serious threat to the model's validity."),
    ("2","Automation bias and algorithm aversion contradict",
     "The literature supports BOTH — people over-trust algorithms and people under-trust them. L1 assumes over-trust. The model currently has no account of when each occurs.",
     "Specify the moderators rather than picking a side: task type, perceived expertise, stakes, and whether the system's output is visible before or after the human forms a view.",
     "Resolving this is a theoretical contribution in its own right."),
    ("3","The sample constraint is structural, not effort",
     "Six eligible per hundred thousand is not fixed by working harder or spending more. A survey design at this prevalence is not viable at any realistic budget.",
     "Two moves, both already underway: a qualitative design for L2 that twelve people CAN sustain, and treating feasibility as a parameter — the GEB 7365 project.",
     "The failure became the research question. That is the strongest thing here.")]
for i,(n,t,prob,fix,note) in enumerate(Wk):
    y=2.32+i*1.48
    rect(s,M,y,SW-2*M,1.38,fill=WHITE,line=RULE,lw=0.7)
    rect(s,M,y,0.05,1.38,fill=GOLD)
    tb(s,M+0.20,y+0.10,0.4,0.3,n,20,GOLD,True)
    tb(s,M+0.55,y+0.09,7,0.28,t,12.5,NAVY,True)
    tb(s,M+0.55,y+0.42,1.6,0.14,"THE PROBLEM",5.8,RED,True,font="Consolas")
    tb(s,M+0.55,y+0.58,4.85,0.50,prob,8,INK)
    tb(s,6.55,y+0.42,1.6,0.14,"WHAT IT NEEDS",5.8,GREEN,True,font="Consolas")
    tb(s,6.55,y+0.58,5.9,0.50,fix,8,INK)
    tb(s,M+0.55,y+1.12,8,0.18,note,7.6,MUTED,italic=True)

# 8 positioning
s=slide(); chrome(s,"WHY THIS RESEARCH, AND WHY NOW","The question the AI-governance market has not answered.","8",
    "Where a doctorate on AI-influenced professional judgment actually sits.")
rect(s,M,2.35,SW-2*M,1.10,fill=NAVY,line=GOLD,lw=1.3)
tb(s,M+0.26,2.52,4,0.2,"THE POSITION",7,GOLD,True,font="Consolas")
tb(s,M+0.26,2.78,11.3,0.6,"Every framework now governing AI — SR 11-7, the NIST AI Risk Management Framework, the EU AI Act — "
   "assumes a competent human reviewer sits above the model. None of them measures whether that reviewer's judgment survives contact with it.",11,RGBColor(0xE3,0xE9,0xF2))
for i,(h,items,col) in enumerate([
    ("WHAT THE FRAMEWORKS ASSUME",["A human reviews model output","Review is independent of the model","Oversight is a control that works"],RED),
    ("WHAT THIS PROGRAMME ASKS",["Does the reviewer defer? (L1)","Does agreement stop the checking? (L2)","Does the evidence base thin? (L3)"],GOLD),
    ("WHY IT IS CREDIBLE HERE",["15 years audit — Citi, JPMorgan","Former state bank examiner","A study that failed honestly and said so"],GREEN)]):
    x=M+i*(3.83+0.12)
    rect(s,x,3.62,3.83,1.60,fill=WHITE,line=col,lw=1.1)
    tb(s,x+0.18,3.78,3.5,0.18,h,6.5,col,True,font="Consolas")
    for j,it in enumerate(items):
        tb(s,x+0.18,4.06+j*0.36,0.2,0.2,"›",8,col,True,font="Consolas")
        tb(s,x+0.38,4.06+j*0.36,3.3,0.32,it,8.6,INK)
rect(s,M,5.42,SW-2*M,1.00,fill=RGBColor(0xF7,0xF4,0xEB),line=GOLD,lw=1.1)
tb(s,M+0.26,5.58,6,0.2,"THE ONE-SENTENCE VERSION",7,GOLD,True,font="Consolas")
tb(s,M+0.26,5.84,11.3,0.5,"Auditing is the profession that already knows how to test whether a control works. This programme "
   "turns that competence on the control everyone is now relying on and nobody has measured: the human being asked to check the machine.",11.5,NAVY,True)

OUT='/home/user/Profile/dba/00_Execution/RESEARCH_STATUS_DECK.pptx'
prs.save(OUT); print("wrote",OUT,len(prs.slides.__iter__.__self__._sldIdLst),"slides")
